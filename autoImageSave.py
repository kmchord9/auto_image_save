import win32clipboard
import win32con
from PIL import ImageGrab
import datetime
import time
import pywintypes
from pptx import Presentation
from pptx.util import Inches, Pt 
from pptx.util import Cm
import os
import re
import keyboard

USERPROFILE = os.environ['USERPROFILE']
SAVE_PATH = f"{USERPROFILE}\\Pictures\\autoImgSave\\"

def pptxAddImage(imgPath, text=None):
    isOpen = False
    now = datetime.datetime.now()
    pptxFileName = now.strftime('%Y%m%d')
    pptxSavePath = f"{SAVE_PATH}{pptxFileName}.pptx"

    if os.path.exists(pptxSavePath):
        prs = Presentation(pptxSavePath)
    else:
        prs = Presentation()

    sld0 = prs.slides.add_slide(prs.slide_layouts[6])

    left = top = width = height = Inches(0.5)
    txBox = sld0.shapes.add_textbox(left, top, width, height) 

    #testBox
    pa = txBox.text_frame.paragraphs[0]

    if text==None:
        pa.text = "This is text inside a textbox"
    else:
        pa.text=text
    pa.font.size = Pt(28)

    pic0 = sld0.shapes.add_picture(imgPath,Cm(1), Cm(3))

    try:
        prs.save(pptxSavePath)

    except PermissionError as e:
        isOpen = True
        return isOpen

    return isOpen

def pptxAddLink(url,title):
    now = datetime.datetime.now()
    pptxFileName = now.strftime('%Y%m%d')
    pptxSavePath = f"{SAVE_PATH}{pptxFileName}.pptx"

    if os.path.exists(pptxSavePath):
        prs = Presentation(pptxSavePath)
    else:
        prs = Presentation()
    
    slide = prs.slides[-1] 
    text_box = slide.shapes.add_textbox(0, 0, 1, 1)
    text_frame = text_box.text_frame
    text_frame.clear()

    # add text
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title

    # add hyperlink
    hyperlink = run.hyperlink
    hyperlink.address = url
    print(title)
    print(url)

    prs.save(pptxSavePath)

def saveResizedImg(img):
    now = datetime.datetime.now()
    fname = now.strftime('%Y%m%d%H%M%S')
    resizedImg = imgResize(img)
    imgPath = f'{SAVE_PATH}{fname}.png'
    resizedImg.save(imgPath)

    return imgPath

def imgResize(img):
    MAX_WIDTH = 680
    MAX_HEIGHT = 450
     
    imgWidth, imgHeight = img.size

    if imgWidth>MAX_WIDTH or imgHeight>MAX_HEIGHT:
        xRatio = MAX_WIDTH/imgWidth
        yRatio = MAX_HEIGHT/imgHeight

        if xRatio < yRatio:
            size = (MAX_WIDTH, round(imgHeight*xRatio))
        else:
            size = (round(imgWidth*yRatio), MAX_HEIGHT)
        resizedImg = img.resize(size)

        return resizedImg
    else:
        return img


def main():
    if not os.path.exists(SAVE_PATH):
        os.makedirs(SAVE_PATH)
    try:
        pptPageTitle=None
        #PPTが開いているときに追加を失敗した画像のパスリスト
        que = []

        #起動時のクリップボード内のデータを取得
        win32clipboard.OpenClipboard()
        if win32clipboard.IsClipboardFormatAvailable(win32con.CF_DIB):
            clip0 = win32clipboard.GetClipboardData(win32con.CF_DIB)
            text0 =""
            
        elif win32clipboard.IsClipboardFormatAvailable(win32con.CF_UNICODETEXT):
            text0 = win32clipboard.GetClipboardData(win32con.CF_UNICODETEXT)
            clip0 =""
        else:
            clip0=""
            text0=""
    finally:
        win32clipboard.CloseClipboard()

    #クリップボード監視開始
    while True:
        if keyboard.is_pressed('shift+escape'):
            pptPageTitle = input("pptのタイトルを入力>> ")     
        try:        
            win32clipboard.OpenClipboard()
            #クリップボード画像の場合
            if win32clipboard.IsClipboardFormatAvailable(win32con.CF_DIB):
                clip1 = win32clipboard.GetClipboardData(win32con.CF_DIB)
                if clip0!=clip1:
                    img = ImageGrab.grabclipboard()
                    imgPath = saveResizedImg(img)                
                    clip0=clip1                   
                    isOpenPPt = pptxAddImage(imgPath,text=pptPageTitle)
                    if isOpenPPt:
                        print("pptが開いています")
                        que.append(imgPath)
                    else:
                        if que:
                            for q in que:
                                isOpenPPt = pptxAddImage(q,text=pptPageTitle)
                            que=[]                
                    print(f"saved:{imgPath}") 
                    continue

            #クリップボードテキストの場合
            elif win32clipboard.IsClipboardFormatAvailable(win32con.CF_UNICODETEXT):
                text1 = win32clipboard.GetClipboardData(win32con.CF_UNICODETEXT)
                if text0!=text1:
                    #マークアップ形式のリンクを抽出                 
                    match = re.search("\[(.*)\]\((.*)\)", text1)
                    if match:
                        title = match.group(1)
                        url = match.group(2)
                        pptxAddLink(url,title)                   
                        text0=text1
                        print(f"link added:{text0}")             
        #クリップボードへのアクセスエラー
        except pywintypes.error as e:
            print(e)
            time.sleep(1)
            continue

        #Ctrl+C時
        except KeyboardInterrupt as e:
            print(e)

        else:
            win32clipboard.CloseClipboard()

        time.sleep(0.5)

if __name__ == "__main__":
    main()